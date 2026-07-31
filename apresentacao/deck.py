# -*- coding: utf-8 -*-
"""
Gera 'Como eu uso IA no dia a dia' sobre o template corporativo NTT DATA.

Usa os layouts do proprio template (master, tema, logo, rodape preservados),
respeitando as regras do slide de uso do template:
  - Georgia em titulos, Arial no corpo
  - nao alterar elementos do master
  - NAO alternar fundo claro e escuro no mesmo deck  -> ESTILO define um so
  - nao cobrir elementos de marca

Uso:  python deck_ntt.py <template.pptx> <saida.pptx> [claro|escuro]
"""
import sys, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE = sys.argv[1]
SAIDA    = sys.argv[2]
ESTILO   = (sys.argv[3] if len(sys.argv) > 3 else "claro").lower()
DARK     = ESTILO == "escuro"

# ---------- paleta da marca (theme1 do template) ----------
NAVY    = "070F26"  # Smart Navy      (background / text)
BRANCO  = "FFFFFF"  # White
AZUL    = "0072BC"  # Future Blue
AZUL_E  = "0055B6"  # Future Blue 160  (contraste em texto)
CIANO   = "00DFED"  # Turquoise
VERDE   = "06B941"  # Green 160
LARANJA = "FF7A00"  # Orange
ERRO    = "B22000"  # Orange 160       (attention grabbing)
CINZA   = "949494"  # Grey 100
CINZA_C = "E8E8E8"  # Grey 60
GRAFITE = "2E404D"  # Text Grey

# tints derivados, para fundo de cartao
if DARK:
    FUNDO, TEXTO, SUAVE = NAVY, BRANCO, "16203A"
    MUTED    = "A8B0BF"
    ERR_BG = OK_BG = NEU_BG = "16203A"
    BORDA_E = BORDA_O = "2E404D"
else:
    FUNDO, TEXTO, SUAVE = BRANCO, NAVY, "F2F5F8"
    MUTED    = "5A6675"
    ERR_BG = OK_BG = NEU_BG = "F2F5F8"
    BORDA_E = BORDA_O = CINZA_C

FT = "Georgia"      # titulos  (regra da marca)
FB = "Arial"        # corpo    (regra da marca)
FM = "Courier New"  # codigo

W, H = 13.333, 7.5
M = 0.72

prs = Presentation(TEMPLATE)

# ---------- mantem a CAPA PADRAO do template (slide 2) e remove o resto ----------
# O template traz 9 slides de exemplo; o 2o e a capa oficial
# ("Full Image, Full Innovation Curve"). Ela e preservada e editada,
# nao recriada, para nao mexer em elemento de marca.
CAPA_IDX = 1
xml_slides = prs.slides._sldIdLst
for i, sld in enumerate(list(xml_slides)):
    if i == CAPA_IDX:
        continue
    rId = sld.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    xml_slides.remove(sld)
CAPA = prs.slides[0]

# A capa preservada ocupa 'slide2.xml'. Os slides novos que o python-pptx cria
# sao numerados a partir de slide1.xml e colidiriam com ela, sobrescrevendo-a
# no pacote. Move a capa para um numero alto antes de adicionar qualquer slide.
from pptx.opc.packuri import PackURI
CAPA.part.partname = PackURI('/ppt/slides/slide900.xml')

# ---------- escolhe a familia de masters conforme o estilo ----------
masters = prs.slide_masters
def lay(mi, nome):
    for l in masters[mi].slide_layouts:
        if l.name.strip().lower().startswith(nome.strip().lower()):
            return l
    raise SystemExit("layout nao encontrado: M%d '%s'" % (mi, nome))

if DARK:
    L_TITULO  = lay(9,  "Title Only")
    L_SUB     = lay(9,  "Title, 1 Column Body Text")
    L_DIV     = lay(13, "Future Blue")   # divisor de secao
    L_FIM     = lay(13, "Smart Navy")    # encerramento
else:
    L_TITULO  = lay(3,  "Title Only")
    L_SUB     = lay(3,  "Title, 1 Column Body Text")
    L_DIV     = lay(13, "Smart Navy")    # divisor de secao
    L_FIM     = lay(13, "Future Blue")   # encerramento (padrao do template)

# ---------- helpers ----------
def rgb(h): return RGBColor.from_string(h)

def add(slide, x, y, w, h):
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))

def txt(slide, s, x, y, w, h, size=14, font=FB, cor=None, bold=False,
        italic=False, align=PP_ALIGN.LEFT, space=None, anchor=MSO_ANCHOR.TOP):
    tb = add(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    linhas = s.split("\n")
    for i, ln in enumerate(linhas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space: p.line_spacing = space
        r = p.add_run(); r.text = ln
        f = r.font
        f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = rgb(cor or TEXTO)
    return tb

def txt_rico(slide, linhas, x, y, w, h, size=11, font=FM, space=1.15):
    """linhas = lista de listas de (trecho, cor). Uma lista interna por linha."""
    tb = add(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, partes in enumerate(linhas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = space
        if isinstance(partes, str):
            partes = [(partes, None)]
        for trecho, cor in partes:
            r = p.add_run(); r.text = trecho
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = rgb(cor or TEXTO)
    return tb


def card(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.04
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    if line:
        sh.line.color.rgb = rgb(line); sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def circulo(slide, x, y, d, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def novo(layout, titulo_txt=None, sub=None):
    s = prs.slides.add_slide(layout)
    # limpa placeholders de corpo que nao vamos usar (evita "Click to edit")
    for ph in list(s.placeholders):
        t = ph.placeholder_format.type
        idx = ph.placeholder_format.idx
        nome = (ph.name or "").lower()
        if str(t).startswith("TITLE") or "title" in nome:
            if titulo_txt:
                ph.text_frame.text = ""
                p = ph.text_frame.paragraphs[0]
                r = p.add_run(); r.text = titulo_txt
                r.font.name = FT; r.font.size = Pt(30); r.font.bold = True
                r.font.color.rgb = rgb(TEXTO)
            else:
                ph._element.getparent().remove(ph._element)
        elif str(t).startswith(("BODY", "OBJECT", "SUBTITLE")):
            ph._element.getparent().remove(ph._element)
    if sub:
        txt(s, sub, M, 1.42, W - 2 * M, 0.4, size=13.5, cor=MUTED, italic=True)
    return s

def bloco_vazio(s, x, y, w, h, marca, titulo, cor, fill, borda):
    """Cartao com cabecalho, sem corpo — o conteudo entra depois com txt_rico."""
    card(s, x, y, w, h, fill, borda)
    circulo(s, x + 0.26, y + 0.2, 0.32, cor)
    txt(s, marca, x + 0.26, y + 0.2, 0.32, 0.32, size=11, bold=True, cor=BRANCO,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, titulo, x + 0.68, y + 0.21, w - 0.94, 0.34, size=14, bold=True, cor=TEXTO)


def bloco(s, x, y, w, h, marca, titulo, corpo, cor, fill, borda, mono=False, size=11.5):
    card(s, x, y, w, h, fill, borda)
    circulo(s, x + 0.26, y + 0.2, 0.32, cor)
    txt(s, marca, x + 0.26, y + 0.2, 0.32, 0.32, size=11, bold=True, cor=BRANCO,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, titulo, x + 0.68, y + 0.21, w - 0.94, 0.34,
        size=14, bold=True, cor=TEXTO)
    txt(s, corpo, x + 0.26, y + 0.66, w - 0.52, h - 0.9,
        size=size, font=(FM if mono else FB), cor=TEXTO, space=1.15 if mono else 1.25)

def rodape(s, t):
    txt(s, t, M, H - 1.02, W - 2.4, 0.38, size=11.5, cor=MUTED, italic=True)

def notas(s, t):
    s.notes_slide.notes_text_frame.text = t

def divisor(letra, titulo_txt, sub):
    # Layout de conteudo, NAO os layouts de marca (Smart Navy / Future Blue):
    # aqueles trazem o logo centralizado e qualquer texto no meio o encobre.
    s = prs.slides.add_slide(L_TITULO)
    for ph in list(s.placeholders):
        if str(ph.placeholder_format.type).startswith(("TITLE", "BODY", "OBJECT", "SUBTITLE")):
            ph._element.getparent().remove(ph._element)
    card(s, M, 2.35, W - 2 * M, 2.3, AZUL)
    circulo(s, M + 0.55, 3.0, 1.0, BRANCO)
    txt(s, letra, M + 0.55, 3.0, 1.0, 1.0, size=32, font=FT, bold=True, cor=AZUL_E,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, titulo_txt, M + 2.0, 3.0, W - M - 3.0, 0.7, size=30, font=FT, bold=True, cor=BRANCO)
    txt(s, sub, M + 2.0, 3.78, W - M - 3.0, 0.5, size=14, cor="D6E9F7", italic=True)
    return s

# =====================================================================
# 1 CAPA  (slide padrao do template, apenas preenchido)
SEP = chr(10)

def preencher_ph(slide, idx, texto, size, font, bold=False, cor=None):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            linhas = texto.split(SEP)
            for i, ln in enumerate(linhas):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                r = p.add_run(); r.text = ln
                r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
                if cor: r.font.color.rgb = rgb(cor)
            return True
    return False

preencher_ph(CAPA, 0,  "Boas práticas no" + SEP + "uso de IA", 34, FT, True, BRANCO)
preencher_ph(CAPA, 11, "Troubleshooting e desenvolvimento de features" + SEP +
                       "com foco em engenharia de dados", 13, FB, False, BRANCO)
notas(CAPA, "Abrir pelo problema do time, nao por credencial pessoal. 30 min de fala, 30 de conversa.")


# ============================ ABERTURA ============================

# 2 POR QUE AGORA — dado de mercado
s = novo(L_TITULO, "Por que essa conversa, agora",
         "DORA 2025 · State of AI-assisted Software Development · Google Cloud")
stats = [("90%", "dos profissionais de\nsoftware já usam IA", AZUL),
         ("2h", "por dia, na mediana,\ntrabalhando com IA", AZUL),
         ("↑", "throughput de entrega\nsubiu com a adoção", VERDE)]
sw, sg = 3.82, 0.33
sx = M
for val, desc, cor in stats:
    card(s, sx, 2.0, sw, 1.75, SUAVE)
    txt(s, val, sx + 0.28, 2.2, sw - 0.56, 0.7, size=34, font=FT, bold=True, cor=cor,
        align=PP_ALIGN.CENTER)
    txt(s, desc, sx + 0.28, 2.95, sw - 0.56, 0.7, size=12, cor=MUTED,
        align=PP_ALIGN.CENTER, space=1.25)
    sx += sw + sg
card(s, M, 4.05, W - 2 * M, 1.15, ERR_BG, BORDA_E)
circulo(s, M + 0.3, 4.35, 0.42, ERRO)
txt(s, "!", M + 0.3, 4.35, 0.42, 0.42, size=15, bold=True, cor=BRANCO,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, "E a estabilidade da entrega caiu.", M + 0.92, 4.28, 10.5, 0.35,
    size=15, bold=True, cor=TEXTO)
txt(s, "Mais volume sem teste automatizado e feedback rápido vira mais incidente.",
    M + 0.92, 4.68, 10.5, 0.35, size=12.5, cor=MUTED)
rodape(s, "Fonte: DORA 2025 — State of AI-assisted Software Development, Google Cloud.")
notas(s, "Abrir pelo dado, nao por opiniao. O numero de estabilidade e o gancho: o problema nao e adotar, e adotar sem base.")

# 3 A TESE
s = novo(L_TITULO, "IA amplifica o que já existe")
txt(s, "O valor não vem da ferramenta. Vem das práticas ao redor dela.",
    M, 1.75, 11.8, 0.5, size=17, cor=MUTED, italic=True)
tw = 5.85
card(s, M, 2.4, tw, 2.15, SUAVE)
txt(s, "Time com base boa", M + 0.3, 2.62, tw - 0.6, 0.35, size=15, bold=True, cor=VERDE)
txt(s, "teste automatizado, versionamento maduro,\nfeedback rápido, escopo escrito\n\n→ entrega mais e mantém estabilidade",
    M + 0.3, 3.05, tw - 0.6, 1.3, size=13, cor=TEXTO, space=1.3)
card(s, M + tw + 0.35, 2.4, tw, 2.15, ERR_BG, BORDA_E)
txt(s, "Time sem base", M + tw + 0.65, 2.62, tw - 0.6, 0.35, size=15, bold=True, cor=ERRO)
txt(s, "sem teste, sem escopo, sem review de verdade\n\n\n→ entrega mais rápido o problema também",
    M + tw + 0.65, 3.05, tw - 0.6, 1.3, size=13, cor=TEXTO, space=1.3)
txt(s, "A pergunta não é “usar ou não usar IA”. É o que precisa existir antes para que ela ajude.",
    M, 4.85, W - 2 * M, 0.5, size=16, bold=True)
notas(s, "Esta e a tese. Tudo depois disso e detalhamento dela.")

# 4 AS DUAS FRENTES
s = novo(L_TITULO, "Duas frentes, dois modos de usar",
         "O que funciona em uma não funciona na outra")
fw = 5.85
card(s, M, 2.0, fw, 2.6, AZUL)
txt(s, "TROUBLESHOOTING", M + 0.32, 2.25, fw - 0.64, 0.35, size=14, bold=True, cor=BRANCO)
txt(s, "O problema já existe.\nVocê precisa achar a causa.\n\nA régua é: quanto você já sabe\nsobre onde está o erro?",
    M + 0.32, 2.75, fw - 0.64, 1.6, size=13.5, cor="CFE6F7", space=1.35)
card(s, M + fw + 0.35, 2.0, fw, 2.6, SUAVE)
txt(s, "FEATURES", M + fw + 0.67, 2.25, fw - 0.64, 0.35, size=14, bold=True, cor=AZUL_E)
txt(s, "O problema ainda não existe.\nVocê precisa construir certo.\n\nA régua é: o projeto tem\ncontexto documentado?",
    M + fw + 0.67, 2.75, fw - 0.64, 1.6, size=13.5, cor=TEXTO, space=1.35)
txt(s, "Misturar as duas é o erro mais comum — e o mais caro.",
    M, 4.85, W - 2 * M, 0.5, size=16, bold=True)
notas(s, "Deixar claro que a agenda segue essas duas partes.")

# ======================= PARTE 1 · TROUBLESHOOTING =======================

notas(divisor("1", "Troubleshooting", "Quando o problema já existe e alguém precisa achar a causa"),
      "Transicao curta.")

# 6 A REGUA
s = novo(L_TITULO, "A régua: quanto você já sabe?",
         "A resposta muda a estratégia — e o custo — completamente")
níveis = [
    ("BAIXA", "Você sabe o arquivo\ne a linha", "Recorte e peça o fix.\nIA é digitação rápida.", VERDE),
    ("MÉDIA", "Você sabe o módulo,\nnão a causa", "Dê o recorte + hipóteses.\nDeixe a IA cruzar.", AZUL),
    ("ALTA", "Você não faz ideia\nde onde começa", "Aqui a IA compensa mais.\nDeixe varrer — com limite.", "B45309"),
]
nw, ng = 3.82, 0.33
nx = M
for nome, quando, oquefazer, cor in níveis:
    card(s, nx, 2.0, nw, 2.6, SUAVE)
    txt(s, nome, nx + 0.26, 2.22, nw - 0.52, 0.35, size=15, bold=True, cor=cor)
    txt(s, quando, nx + 0.26, 2.68, nw - 0.52, 0.7, size=12.5, cor=MUTED, space=1.25)
    txt(s, oquefazer, nx + 0.26, 3.5, nw - 0.52, 0.9, size=13, cor=TEXTO, space=1.3)
    nx += nw + ng
txt(s, "A régua é pessoal: o que é alta complexidade para um é média para outro.",
    M, 4.85, W - 2 * M, 0.45, size=15, bold=True)
txt(s, "Por isso ninguém pode te dar o número pronto. O que dá para dar é o critério.",
    M, 5.3, 11.8, 0.45, size=13, cor=MUTED, italic=True)
notas(s, "Ponto do Giovanni: se voce nao tem nocao do que esta acontecendo, use IA para varrer. Se tem, va direto e economize.")

# 7 O CUSTO DE NAO RECORTAR
s = novo(L_TITULO, "Complexidade baixa: o custo de não recortar")
TS, LOG_I, LOG_E, CODE, PROSA = CINZA, AZUL, ERRO, "6A3FA0", TEXTO
bloco_vazio(s, M, 1.72, 5.8, 2.85, "✕", "Log inteiro no prompt", ERRO, ERR_BG, BORDA_E)
txt_rico(s, [
    [("Aqui o log, me ajuda:", PROSA)],
    [("", None)],
    [("09:14:02 ", TS), ("INFO ", LOG_I), (" Starting worker", PROSA)],
    [("09:14:02 ", TS), ("INFO ", LOG_I), (" Connected redis", PROSA)],
    [("  ... (1.847 linhas) ...", TS)],
    [("09:31:55 ", TS), ("ERROR ", LOG_E), ("ValueError", PROSA)],
], M + 0.26, 2.38, 5.3, 2.0, size=10.5)
bloco_vazio(s, M + 6.1, 1.72, 5.8, 2.85, "✓", "Recorte + o que já verificou", VERDE, OK_BG, BORDA_O)
txt_rico(s, [
    [("Erro no worker, ~2% dos registros.", PROSA)],
    [("", None)],
    [("ValueError", LOG_E), (": invalid literal", PROSA)],
    [("  parser.py:42", CODE), (" in ", PROSA), ("parse_row", CODE)],
    [("", None)],
    [("Já verifiquei: CSV tem célula vazia.", PROSA)],
    [("Schema é do fornecedor, não muda.", PROSA)],
], M + 6.36, 2.38, 5.3, 2.0, size=10.5)
txt(s, "~40.000 tokens", M, 4.72, 5.8, 0.4, size=17, bold=True, cor=ERRO, align=PP_ALIGN.CENTER)
txt(s, "~200 tokens", M + 6.1, 4.72, 5.8, 0.4, size=17, bold=True, cor=AZUL_E, align=PP_ALIGN.CENTER)
txt(s, "As 1.847 linhas irrelevantes disputam atenção com as 6 que importam. Contexto cheio de ruído piora a resposta antes de encher a janela.",
    M, 5.25, 11.9, 0.6, size=13.5, cor=MUTED, space=1.3)
notas(s, "Aqui o ponto nao e economia: e qualidade da resposta.")

# 8 COMPLEXIDADE ALTA
s = novo(L_TITULO, "Complexidade alta: aqui a IA compensa",
         "Quando você não sabe por onde começar, varrer é justamente o que ela faz melhor")
cw2 = 5.85
card(s, M, 2.0, cw2, 2.5, SUAVE)
txt(s, "Na mão", M + 0.3, 2.22, cw2 - 0.6, 0.35, size=15, bold=True, cor=MUTED)
txt(s, "3 a 6 horas caçando\nem repositório que ninguém domina\n\nCusto: hora de engenheiro sênior",
    M + 0.3, 2.68, cw2 - 0.6, 1.5, size=13.5, cor=TEXTO, space=1.35)
card(s, M + cw2 + 0.35, 2.0, cw2, 2.5, AZUL)
txt(s, "Com agente", M + cw2 + 0.65, 2.22, cw2 - 0.6, 0.35, size=15, bold=True, cor=BRANCO)
txt(s, "Varre, cruza chamadas, propõe hipótese\n\n≈ 12 ACU numa sessão de 3h\n≈ US$ 25",
    M + cw2 + 0.65, 2.68, cw2 - 0.6, 1.5, size=13.5, cor="CFE6F7", space=1.35)
txt(s, "Quando a alternativa é meio dia de investigação, gastar token é a decisão barata.",
    M, 4.75, W - 2 * M, 0.45, size=16, bold=True)
txt(s, "Não existe cálculo exato — existe ordem de grandeza. E ela costuma favorecer a IA quanto menos você souber.",
    M, 5.2, 11.9, 0.45, size=13, cor=MUTED, italic=True)
rodape(s, "1 ACU ≈ 15 min de trabalho ativo do agente · referência de preço pública do fornecedor.")
notas(s, "Ponto da Camila: mesmo quando voce tem nocao, pode levar tempo demais. Comparar com hora humana, nao com zero.")

# 9 CASO DE DADOS
s = novo(L_TITULO, "Caso real: pipeline com número divergente")
card(s, M, 1.78, W - 2 * M, 0.92, SUAVE)
txt(s, "“O faturamento de junho no dashboard está R$ 200 acima do fechamento contábil.”",
    M + 0.32, 2.02, 11.4, 0.45, size=15, italic=True)
passos2 = [
    ("1", "Recortar", "Só junho, só as colunas do agregado.\nNão o pipeline inteiro."),
    ("2", "Dar o que já sabe", "Houve reprocessamento de lote em 14/06.\nA IA não tinha como adivinhar."),
    ("3", "Pedir hipóteses", "Antes do fix: liste as causas possíveis\ne como distinguir uma da outra."),
]
py = 3.0
for num, tit, desc in passos2:
    circulo(s, M, py, 0.42, AZUL)
    txt(s, num, M, py, 0.42, 0.42, size=12, bold=True, cor=BRANCO,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, tit, M + 0.65, py - 0.02, 2.6, 0.35, size=14, bold=True)
    txt(s, desc, M + 3.4, py - 0.04, 8.2, 0.7, size=12.5, cor=MUTED, space=1.25)
    py += 0.85
txt(s, "A causa era duplicata de reprocessamento. O agente achou em minutos — porque recebeu o recorte e a pista.",
    M, 5.65, 11.9, 0.5, size=14, italic=True)
notas(s, "Exemplo de dados, nao de backend. Foi pedido explicito da rodada anterior.")

# ========================= PARTE 2 · FEATURES =========================

notas(divisor("2", "Features", "Quando o problema ainda não existe e você precisa construir certo"),
      "Transicao curta.")

# 11 GREENFIELD x BROWNFIELD
s = novo(L_TITULO, "Greenfield e brownfield pedem coisas diferentes",
         "Greenfield: projeto do zero · Brownfield: sistema que já existe e você altera")
gw2 = 5.85
card(s, M, 2.15, gw2, 2.5, SUAVE)
txt(s, "GREENFIELD", M + 0.3, 2.38, gw2 - 0.6, 0.35, size=14, bold=True, cor=VERDE)
txt(s, "Sem legado, sem armadilha herdada.\n\nA IA rende muito desde o primeiro dia —\ne é aqui que o SDD brilha.\n\nRisco: velocidade sem decisão de arquitetura.",
    M + 0.3, 2.85, gw2 - 0.6, 1.6, size=13, cor=TEXTO, space=1.3)
card(s, M + gw2 + 0.35, 2.15, gw2, 2.5, ERR_BG, BORDA_E)
txt(s, "BROWNFIELD", M + gw2 + 0.65, 2.38, gw2 - 0.6, 0.35, size=14, bold=True, cor=ERRO)
txt(s, "Anos de decisão não documentada.\n\nPedir feature direto aqui é o erro clássico:\no agente lê o código e inventa o resto.\n\nPrecisa de um passo antes.",
    M + gw2 + 0.65, 2.85, gw2 - 0.6, 1.6, size=13, cor=TEXTO, space=1.3)
txt(s, "A maior parte do trabalho real é brownfield. E é onde quase ninguém faz o passo anterior.",
    M, 4.95, W - 2 * M, 0.5, size=16, bold=True)
notas(s, "Definir os dois termos em voz alta: nem todo mundo conhece.")

# 12 BROWNFIELD: DOCUMENTAR ANTES
s = novo(L_TITULO, "Brownfield: documentar antes de pedir feature",
         "Uma vez por projeto. Depois disso toda tarefa fica mais barata")
docs = [
    ("Contexto do projeto", "AGENTS.md na raiz: como rodar, como testar,\narmadilhas, convenções não óbvias, glossário", AZUL),
    ("Conhecimento de negócio", "O que o código não conta: regra que veio do\njurídico, decisão antiga e o porquê dela", AZUL),
    ("Ambiente reproduzível", "Setup do repositório e snapshot de máquina,\npara o agente não gastar sessão descobrindo build", AZUL),
]
dy = 2.15
for tit, desc, cor in docs:
    card(s, M, dy, W - 2 * M, 0.95, SUAVE)
    txt(s, tit, M + 0.3, dy + 0.18, 3.6, 0.35, size=14, bold=True, cor=cor)
    txt(s, desc, M + 4.1, dy + 0.14, 7.5, 0.65, size=12.5, cor=MUTED, space=1.25)
    dy += 1.08
txt(s, "O agente pode escrever boa parte disso — lendo o repositório e te perguntando o que não dá para inferir.",
    M, 5.5, 11.9, 0.5, size=14, italic=True)
notas(s, "AGENTS.md e padrao aberto e o Devin le nativamente. Vale citar.")

# 13 SDD — O QUE E
s = novo(L_TITULO, "SDD — Spec-Driven Development",
         "A especificação é a fonte de verdade. O código é saída dela")
txt(s, "Quando a IA escreve o código, escrever código deixa de ser o gargalo.",
    M, 2.0, 11.8, 0.45, size=19, font=FT)
txt(s, "Decidir o que deve ser construído vira o gargalo.",
    M, 2.55, 11.8, 0.45, size=19, font=FT, bold=True, cor=AZUL)
card(s, M, 3.3, W - 2 * M, 1.5, SUAVE)
txt(s, "Sem spec, quem decide é o modelo — por omissão, com confiança, e você descobre no review.",
    M + 0.32, 3.55, 11.4, 0.4, size=14.5, bold=True)
txt(s, "Não é documentação a mais. É a conversa que já precisaria acontecer, escrita uma vez\ne reaproveitada por todo mundo — inclusive pelo próximo agente.",
    M + 0.32, 4.0, 11.4, 0.7, size=13, cor=MUTED, space=1.3)
notas(s, "Dois ou tres slides no maximo sobre SDD. Este e o conceito.")

# 14 SDD — O CICLO
s = novo(L_TITULO, "O ciclo do SDD",
         "O mesmo, de tarefa pequena a feature grande — só muda o rigor de cada etapa")
import math
TONS = ["0055B6", "0072BC", "19A3FC", "00C3E0", "7BF7FF"]
CX, CY, RAIO, D = 4.05, 3.75, 1.62, 1.34
passos = [("ENTENDER", "ler antes de pedir"), ("ESPECIFICAR", "escrever o critério"),
          ("EXECUTAR", "pessoa ou agente"), ("VERIFICAR", "rodar, não ler"),
          ("REGISTRAR", "o que vale pra próxima")]
N = len(passos); PASSO = 360.0 / N
for i in range(N):
    ini = -90 - PASSO / 2 + i * PASSO
    arco = s.shapes.add_shape(MSO_SHAPE.BLOCK_ARC,
                              Inches(CX - RAIO - 0.34), Inches(CY - RAIO - 0.34),
                              Inches((RAIO + 0.34) * 2), Inches((RAIO + 0.34) * 2))
    arco.adjustments[0] = ini * 0.6
    arco.adjustments[1] = (ini + PASSO - 3) * 0.6
    arco.adjustments[2] = 0.19
    arco.fill.solid(); arco.fill.fore_color.rgb = rgb(TONS[i])
    arco.line.fill.background(); arco.shadow.inherit = False
for i, (nome, desc) in enumerate(passos):
    ang = math.radians(-90 + i * PASSO)
    x = CX + RAIO * math.cos(ang) - D / 2
    y = CY + RAIO * math.sin(ang) - D / 2
    dest = (i == 1)
    circulo(s, x, y, D, AZUL if dest else BRANCO)
    aro = s.shapes[-1]; aro.line.color.rgb = rgb(TONS[i]); aro.line.width = Pt(2.5)
    txt(s, nome, x + 0.04, y + 0.3, D - 0.08, 0.34, size=8.5, bold=True,
        cor=BRANCO if dest else GRAFITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, str(i + 1), x + 0.04, y + 0.62, D - 0.08, 0.28, size=9,
        cor="CFE6F7" if dest else CINZA, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, "repete\na cada\ntarefa", CX - 0.75, CY - 0.42, 1.5, 0.9, size=11, font=FT,
    italic=True, cor=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=1.15)
ly = 2.5
for i, (nome, desc) in enumerate(passos):
    dest = (i == 1)
    circulo(s, 7.35, ly + 0.02, 0.3, TONS[i])
    txt(s, str(i + 1), 7.35, ly + 0.02, 0.3, 0.3, size=10, bold=True,
        cor=BRANCO if i < 3 else GRAFITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, nome, 7.82, ly, 2.3, 0.3, size=12.5, bold=True, cor=AZUL_E if dest else TEXTO)
    txt(s, desc, 10.1, ly + 0.01, 2.6, 0.3, size=11.5, cor=MUTED)
    ly += 0.58
txt(s, "A etapa 2 é a que mais se pula — e a única que, quando falha, estraga todas as outras.",
    M, 5.72, W - 2 * M, 0.5, size=15, italic=True)
notas(s, "Mostrar que em tarefa pequena cada etapa leva segundos. Nao e cerimonia.")

# 15 FRAMEWORKS DE SDD
s = novo(L_TITULO, "Frameworks que implementam o ciclo",
         "Impõem o método dentro da ferramenta, em vez de depender de disciplina")
fws = [("Superpowers", "Skills que conduzem brainstorming,\nplano, TDD e review.\nEscreve a spec sozinho."),
       ("Spec Kit", "Fluxo em fases explícitas:\nespecificar, planejar, tarefas.\nCLI própria."),
       ("GSD", "Mais enxuto e opinativo.\nFoco em fechar escopo\nantes de gerar código.")]
fw2, fg = 3.82, 0.33
fx = M
for nome, desc in fws:
    card(s, fx, 2.15, fw2, 2.2, SUAVE)
    txt(s, nome, fx + 0.28, 2.38, fw2 - 0.56, 0.35, size=15, bold=True, cor=AZUL_E)
    txt(s, desc, fx + 0.28, 2.85, fw2 - 0.56, 1.3, size=12.5, cor=MUTED, space=1.3)
    fx += fw2 + fg
txt(s, "Escolha um e use por um mês antes de comparar. O ganho vem da consistência, não do nome.",
    M, 4.65, W - 2 * M, 0.5, size=15, bold=True)
txt(s, "Todos cobrem brainstorming, spec e plano. Nenhum cobre delegação a agente externo nem memória entre sessões.",
    M, 5.15, 11.9, 0.5, size=13, cor=MUTED, italic=True)
notas(s, "Nao vender um. Mostrar que existe categoria e que a escolha importa menos que a consistencia.")

# 16 ANATOMIA DE UMA BOA SPEC
s = novo(L_TITULO, "Anatomia de uma especificação que funciona")
anat = [("Problema", "com evidência: número, frequência, impacto"),
        ("Comportamento esperado", "o que passa a ser verdade quando estiver pronto"),
        ("Regras", "as decisões que o modelo não tem como adivinhar"),
        ("Fora de escopo", "o que não pode ser tocado — a seção que mais economiza"),
        ("Aceite", "verificável por comando, não por opinião")]
ay = 1.95
for i, (tit, desc) in enumerate(anat):
    dest = (i == 3)
    card(s, M, ay, W - 2 * M, 0.72, ERR_BG if dest else SUAVE, BORDA_E if dest else None)
    txt(s, tit, M + 0.3, ay + 0.19, 3.5, 0.35, size=13.5, bold=True,
        cor=ERRO if dest else TEXTO)
    txt(s, desc, M + 4.0, ay + 0.19, 7.6, 0.35, size=12.5, cor=MUTED)
    ay += 0.82
txt(s, "“Fora de escopo” é a primeira a ficar vazia e a que mais evita retrabalho.",
    M, 5.75, W - 2 * M, 0.45, size=15, bold=True)
notas(s, "Se a spec sai sem Fora de escopo, o brainstorming nao terminou.")

# 17 SPEC — EXEMPLO DE DADOS
s = novo(L_TITULO, "O mesmo pedido, com e sem especificação",
         "Deduplicação em pipeline: onde a chave técnica e a chave de negócio divergem")
bloco_vazio(s, M, 1.9, 5.8, 2.75, "✕", "Sem spec", ERRO, ERR_BG, BORDA_E)
txt_rico(s, [
    [("“deduplica a tabela de pedidos”", TEXTO)],
    [("", None)],
    [("PARTITION BY pedido_id", "6A3FA0")],
    [("", None)],
    [("Corrige o número do dashboard.", MUTED)],
    [("Apaga 5 anos de trilha de", ERRO)],
    [("auditoria — e passa no CI.", ERRO)],
], M + 0.26, 2.56, 5.3, 1.9, size=11, font=FB)
bloco_vazio(s, M + 6.1, 1.9, 5.8, 2.75, "✓", "Com spec", VERDE, OK_BG, BORDA_O)
txt_rico(s, [
    [("Regra: duplicata é a repetição de", TEXTO)],
    [("(pedido_id, status, atualizado_em).", TEXTO)],
    [("Fora de escopo: DELETE em bronze.", TEXTO)],
    [("", None)],
    [("PARTITION BY pedido_id,", "6A3FA0")],
    [("  status, atualizado_em", "6A3FA0")],
], M + 6.36, 2.56, 5.3, 1.9, size=11, font=FB)
txt(s, "A diferença cabe numa linha de código — e vem de uma regra que não está em lugar nenhum do repositório.",
    M, 4.85, 11.9, 0.5, size=15, bold=True)
txt(s, "Append-only por requisito de auditoria, decidido com compliance. Nenhum parser de AST descobre isso.",
    M, 5.35, 11.9, 0.45, size=13, cor=MUTED, italic=True)
notas(s, "Exemplo de dados. O erro e previsivel: qualquer agente sem a regra deduplica por pedido_id.")

# 18 TRABALHO CONJUNTO
s = novo(L_TITULO, "Quando a IA não é a única mexendo no código",
         "SDD assume que o agente é o único autor. Na vida real, não é")
card(s, M, 2.05, W - 2 * M, 1.05, ERR_BG, BORDA_E)
txt(s, "O problema", M + 0.3, 2.25, 3.0, 0.32, size=13.5, bold=True, cor=ERRO)
txt(s, "Você altera na mão, o agente não sabe. Na próxima sessão ele relê o repositório inteiro\npara descobrir o que mudou — ou pior, assume o estado antigo.",
    M + 0.3, 2.6, 11.4, 0.6, size=12.5, cor=TEXTO, space=1.25)
card(s, M, 3.3, W - 2 * M, 1.75, SUAVE)
txt(s, "A solução: changelog para agente", M + 0.3, 3.52, 6.0, 0.32, size=13.5, bold=True, cor=AZUL_E)
txt(s, "Um arquivo curto, append-only, com o que mudou fora das sessões de IA:\no que você alterou, por quê, e o que isso invalida do que já foi decidido.\n\nÉ mais barato manter três linhas por mudança do que pagar releitura de repositório toda sessão.",
    M + 0.3, 3.9, 11.4, 1.0, size=12.5, cor=TEXTO, space=1.3)
txt(s, "Vale o mesmo princípio da memória: registre o que ficou verdadeiro, não o que aconteceu.",
    M, 5.25, 11.9, 0.45, size=14, italic=True)
notas(s, "Pedido explicito da rodada anterior. O ponto e nao obrigar a IA a reler o codigo inteiro.")

# ======================= PARTE 3 · DECISAO =======================

notas(divisor("3", "Decidir e verificar", "O que delegar, quanto custa e como revisar o que voltou"),
      "Transicao curta.")

# 20 MATRIZ
s = novo(L_TITULO, "Na mão · piloto · delega",
         "Decisão que vem antes de abrir a ferramenta, não depois")
cols = [("NA MÃO", ERRO, ERR_BG, "Decisão de arquitetura\nRegra de negócio nova\nCódigo de auth e pagamento\nIncidente em produção"),
        ("PILOTO", AZUL, NEU_BG, "Causa desconhecida\nDesign ainda aberto\nRefactor que atravessa módulos\nCódigo sensível, passo a passo"),
        ("DELEGA", VERDE, OK_BG, "Migração mecânica\nBug com repro determinístico\nTeste faltando\nPadronização em N arquivos")]
bw2, gap2 = 3.82, 0.33
x2 = M
for nome, cor, bg, corpo in cols:
    card(s, x2, 2.05, bw2, 2.45, bg)
    txt(s, nome, x2 + 0.24, 2.24, bw2 - 0.48, 0.35, size=14, bold=True, cor=cor)
    txt(s, corpo, x2 + 0.24, 2.7, bw2 - 0.48, 1.7, size=12, space=1.4)
    x2 += bw2 + gap2
card(s, M, 4.7, W - 2 * M, 0.85, SUAVE)
txt(s, "Regra prática dos fornecedores de agente: se você faria em até 3 horas, é delegável. Acima disso, decomponha.",
    M + 0.3, 4.92, 11.4, 0.4, size=14, bold=True)
txt(s, "E a pergunta que resolve quase tudo: existe teste que prova que ficou certo? Se não existe, não delegue.",
    M, 5.65, 11.9, 0.45, size=13, cor=MUTED, italic=True)
notas(s, "Decisao de arquitetura e regra de negocio nunca vao para a IA - ponto levantado na rodada anterior.")

# 21 REVIEW
s = novo(L_TITULO, "Um PR que passa no CI e não deve ser aprovado")
card(s, M, 1.72, 6.5, 2.9, NAVY if not DARK else "16203A")
DIFF_CTX, DIFF_DEL, DIFF_ADD = "C7D0DC", "FF6B4A", "38F990"
txt_rico(s, [
    [("  def parse_row(row) -> int | None:", DIFF_CTX)],
    [("-     return int(row[\"quantidade\"])", DIFF_DEL)],
    [("+     try:", DIFF_ADD)],
    [("+         return int(row[\"quantidade\"])", DIFF_ADD)],
    [("+     except Exception:", DIFF_ADD)],
    [("+         return None", DIFF_ADD)],
    [("", None)],
    [("- def test_invalido():", DIFF_DEL)],
    [("-     with pytest.raises(ValueError):", DIFF_DEL)],
    [("+ def test_invalido():", DIFF_ADD)],
    [("+     assert parse_row(\"abc\") is None", DIFF_ADD)],
], M + 0.28, 1.95, 5.9, 2.45, size=10, space=1.15)
txt(s, "CI verde. Três problemas:", M + 7.0, 1.76, 4.9, 0.4, size=15, bold=True)
txt_rico(s, [
    [("•  ", CINZA), ("except Exception", ERRO), (" engole tudo — coluna", MUTED)],
    [("   faltando vira None silencioso", MUTED)],
    [("", None)],
    [("•  ", CINZA), ("o teste antigo foi invertido", ERRO), (" para passar", MUTED)],
    [("", None)],
    [("•  ", CINZA), ("parâmetro de biblioteca que", MUTED)],
    [("   simplesmente não existe", MUTED)],
], M + 7.0, 2.3, 4.9, 2.2, size=12, font=FB, space=1.25)
txt(s, "Ler o diff não é revisar. O segundo item é o mais grave e o mais fácil de passar batido.",
    M, 4.8, W - 2 * M, 0.45, size=15, bold=True)
txt(s, "Atalho: git diff nos testes primeiro. Linha alterada em teste existente exige justificativa escrita.",
    M, 5.26, 11.8, 0.45, size=13, cor=MUTED, italic=True)
notas(s, "Conduzir como pergunta: mostrar e deixar a plateia achar antes de revelar.")

# 22 BONUS — ESCOLHER MODELO
s = novo(L_TITULO, "Bônus: escolher o modelo pela tarefa",
         "Rodar tudo no topo de linha é o desperdício mais comum e mais fácil de corrigir")
mods = [("Pequeno / rápido", "classificar · extrair · formatar\nresumir · traduzir", VERDE),
        ("Intermediário", "código do dia a dia\nrevisão · refactor local", AZUL),
        ("Topo de linha", "arquitetura · debug difícil\nraciocínio longo · tarefa agêntica", "B45309")]
mw, mg = 3.82, 0.33
mx = M
for nome, uso, cor in mods:
    card(s, mx, 2.15, mw, 1.95, SUAVE)
    txt(s, nome, mx + 0.28, 2.38, mw - 0.56, 0.35, size=14, bold=True, cor=cor)
    txt(s, uso, mx + 0.28, 2.85, mw - 0.56, 1.0, size=12.5, cor=MUTED, space=1.3)
    mx += mw + mg
txt(s, "A diferença de preço entre o menor e o maior costuma ser de uma ordem de grandeza.",
    M, 4.4, W - 2 * M, 0.45, size=15, bold=True)
txt(s, "E o maior nem sempre entrega melhor: em tarefa simples você paga mais, espera mais e recebe o mesmo.",
    M, 4.88, 11.9, 0.45, size=13, cor=MUTED, italic=True)
card(s, M, 5.4, W - 2 * M, 0.72, NEU_BG)
txt(s, "Contexto acima de ~50% da janela aumenta a chance de erro. Sessão nova é mais barata que sessão longa.",
    M + 0.3, 5.58, 11.4, 0.4, size=13, bold=True, cor=AZUL_E)
notas(s, "Slide bonus. Passar rapido se o tempo apertar.")

# 23 POR ONDE COMECAR
s = novo(L_TITULO, "Por onde começar",
         "Nessa ordem — a primeira sozinha já muda o resultado")
inicio = [("1", "Recorte antes de pedir", "No troubleshooting: mande o trecho e o que já verificou."),
          ("2", "Documente o projeto uma vez", "Em brownfield, AGENTS.md antes de qualquer feature."),
          ("3", "Escreva o critério de aceite", "Uma linha verificável por comando, junto com o pedido."),
          ("4", "Rode antes de aprovar", "Ler o diff não é revisar. Vale para PR de IA e de gente.")]
iy = 2.05
for num, tit, desc in inicio:
    card(s, M, iy, W - 2 * M, 0.9, SUAVE)
    circulo(s, M + 0.3, iy + 0.24, 0.44, AZUL)
    txt(s, num, M + 0.3, iy + 0.24, 0.44, 0.44, size=13, bold=True, cor=BRANCO,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, tit, M + 1.0, iy + 0.18, 4.6, 0.35, size=14.5, bold=True)
    txt(s, desc, M + 5.8, iy + 0.2, 5.8, 0.35, size=12.5, cor=MUTED)
    iy += 1.0
txt(s, "Nenhuma das quatro depende de instalar nada.",
    M, 6.05, W - 2 * M, 0.45, size=15, bold=True)
notas(s, "Pedido concreto: escolher UMA e testar essa semana.")

# 24 Q&A
s = novo(L_TITULO, "Agora é de vocês")
txt(s, "30 minutos. Discordar é o melhor uso desse tempo.", M, 1.5, 11, 0.5,
    size=16, bold=True, cor=AZUL_E)
qs = [("Onde isso não funciona no contexto de vocês?", ""),
      ("Qual a régua de complexidade de cada um aqui?", "onde vocês já sentem que vale gastar token"),
      ("Qual tarefa da sprint serviria de teste?", "a mais concreta — começa por essa se travar")]
qy = 2.35
for pergunta, nota in qs:
    card(s, M, qy, W - 2 * M, 0.95, SUAVE)
    txt(s, pergunta, M + 0.32, qy + (0.3 if not nota else 0.16), 10.8, 0.4, size=15, bold=True)
    if nota:
        txt(s, nota, M + 0.32, qy + 0.56, 10.8, 0.3, size=12, cor=MUTED, italic=True)
    qy += 1.12
txt(s, "Material, templates e o exercício prático: github.com/devdiasbr/boas-praticas-ia",
    M, 5.95, 11, 0.4, size=12.5, cor=MUTED, italic=True)
notas(s, "Se travar, comecar pela terceira pergunta.")

# 25 ENCERRAMENTO NTT
s = prs.slides.add_slide(L_FIM)
for ph in list(s.placeholders):
    if str(ph.placeholder_format.type).startswith(("TITLE", "BODY", "OBJECT", "SUBTITLE")):
        ph._element.getparent().remove(ph._element)
notas(s, "Slide de encerramento da marca. Nada sobreposto ao logo.")

prs.save(SAIDA)
print("ok: %s  (%d slides, estilo=%s)" % (SAIDA, len(prs.slides._sldIdLst), ESTILO))
