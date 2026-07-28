# -*- coding: utf-8 -*-
"""
Gera 'Como eu uso IA no dia a dia' sobre o template corporativo NTT DATA.

Usa os layouts do proprio template (master, tema, logo, rodape preservados),
respeitando as regras do slide de uso do template:
  - Georgia em titulos, Arial no corpo
  - nao alterar elementos do master
  - NAO alternar fundo claro e escuro no mesmo deck  -> ESTILO define um so
  - nao cobrir elementos de marca

Uso:  python deck_ntt.py <template.pptx> <saida.pptx> [claro|escuro]
"""
import sys, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE = sys.argv[1]
SAIDA    = sys.argv[2]
ESTILO   = (sys.argv[3] if len(sys.argv) > 3 else "claro").lower()
DARK     = ESTILO == "escuro"

# ---------- paleta da marca (theme1 do template) ----------
NAVY    = "070F26"  # Smart Navy      (background / text)
BRANCO  = "FFFFFF"  # White
AZUL    = "0072BC"  # Future Blue
AZUL_E  = "0055B6"  # Future Blue 160  (contraste em texto)
CIANO   = "00DFED"  # Turquoise
VERDE   = "06B941"  # Green 160
LARANJA = "FF7A00"  # Orange
ERRO    = "B22000"  # Orange 160       (attention grabbing)
CINZA   = "949494"  # Grey 100
CINZA_C = "E8E8E8"  # Grey 60
GRAFITE = "2E404D"  # Text Grey

# tints derivados, para fundo de cartao
if DARK:
    FUNDO, TEXTO, SUAVE = NAVY, BRANCO, "16203A"
    MUTED    = "A8B0BF"
    ERR_BG = OK_BG = NEU_BG = "16203A"
    BORDA_E = BORDA_O = "2E404D"
else:
    FUNDO, TEXTO, SUAVE = BRANCO, NAVY, "F2F5F8"
    MUTED    = "5A6675"
    ERR_BG = OK_BG = NEU_BG = "F2F5F8"
    BORDA_E = BORDA_O = CINZA_C

FT = "Georgia"      # titulos  (regra da marca)
FB = "Arial"        # corpo    (regra da marca)
FM = "Courier New"  # codigo

W, H = 13.333, 7.5
M = 0.72

prs = Presentation(TEMPLATE)

# ---------- mantem a CAPA PADRAO do template (slide 2) e remove o resto ----------
# O template traz 9 slides de exemplo; o 2o e a capa oficial
# ("Full Image, Full Innovation Curve"). Ela e preservada e editada,
# nao recriada, para nao mexer em elemento de marca.
CAPA_IDX = 1
xml_slides = prs.slides._sldIdLst
for i, sld in enumerate(list(xml_slides)):
    if i == CAPA_IDX:
        continue
    rId = sld.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    xml_slides.remove(sld)
CAPA = prs.slides[0]

# A capa preservada ocupa 'slide2.xml'. Os slides novos que o python-pptx cria
# sao numerados a partir de slide1.xml e colidiriam com ela, sobrescrevendo-a
# no pacote. Move a capa para um numero alto antes de adicionar qualquer slide.
from pptx.opc.packuri import PackURI
CAPA.part.partname = PackURI('/ppt/slides/slide900.xml')

# ---------- escolhe a familia de masters conforme o estilo ----------
masters = prs.slide_masters
def lay(mi, nome):
    for l in masters[mi].slide_layouts:
        if l.name.strip().lower().startswith(nome.strip().lower()):
            return l
    raise SystemExit("layout nao encontrado: M%d '%s'" % (mi, nome))

if DARK:
    L_TITULO  = lay(9,  "Title Only")
    L_SUB     = lay(9,  "Title, 1 Column Body Text")
    L_DIV     = lay(13, "Future Blue")   # divisor de secao
    L_FIM     = lay(13, "Smart Navy")    # encerramento
else:
    L_TITULO  = lay(3,  "Title Only")
    L_SUB     = lay(3,  "Title, 1 Column Body Text")
    L_DIV     = lay(13, "Smart Navy")    # divisor de secao
    L_FIM     = lay(13, "Future Blue")   # encerramento (padrao do template)

# ---------- helpers ----------
def rgb(h): return RGBColor.from_string(h)

def add(slide, x, y, w, h):
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))

def txt(slide, s, x, y, w, h, size=14, font=FB, cor=None, bold=False,
        italic=False, align=PP_ALIGN.LEFT, space=None, anchor=MSO_ANCHOR.TOP):
    tb = add(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    linhas = s.split("\n")
    for i, ln in enumerate(linhas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space: p.line_spacing = space
        r = p.add_run(); r.text = ln
        f = r.font
        f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = rgb(cor or TEXTO)
    return tb

def card(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.04
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    if line:
        sh.line.color.rgb = rgb(line); sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def circulo(slide, x, y, d, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def novo(layout, titulo_txt=None, sub=None):
    s = prs.slides.add_slide(layout)
    # limpa placeholders de corpo que nao vamos usar (evita "Click to edit")
    for ph in list(s.placeholders):
        t = ph.placeholder_format.type
        idx = ph.placeholder_format.idx
        nome = (ph.name or "").lower()
        if str(t).startswith("TITLE") or "title" in nome:
            if titulo_txt:
                ph.text_frame.text = ""
                p = ph.text_frame.paragraphs[0]
                r = p.add_run(); r.text = titulo_txt
                r.font.name = FT; r.font.size = Pt(30); r.font.bold = True
                r.font.color.rgb = rgb(TEXTO)
            else:
                ph._element.getparent().remove(ph._element)
        elif str(t).startswith(("BODY", "OBJECT", "SUBTITLE")):
            ph._element.getparent().remove(ph._element)
    if sub:
        txt(s, sub, M, 1.42, W - 2 * M, 0.4, size=13.5, cor=MUTED, italic=True)
    return s

def bloco(s, x, y, w, h, marca, titulo, corpo, cor, fill, borda, mono=False, size=11.5):
    card(s, x, y, w, h, fill, borda)
    circulo(s, x + 0.26, y + 0.2, 0.32, cor)
    txt(s, marca, x + 0.26, y + 0.2, 0.32, 0.32, size=11, bold=True, cor=BRANCO,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, titulo, x + 0.68, y + 0.21, w - 0.94, 0.34,
        size=14, bold=True, cor=TEXTO)
    txt(s, corpo, x + 0.26, y + 0.66, w - 0.52, h - 0.9,
        size=size, font=(FM if mono else FB), cor=TEXTO, space=1.15 if mono else 1.25)

def rodape(s, t):
    txt(s, t, M, H - 1.02, W - 2.4, 0.38, size=11.5, cor=MUTED, italic=True)

def notas(s, t):
    s.notes_slide.notes_text_frame.text = t

def divisor(letra, titulo_txt, sub):
    s = prs.slides.add_slide(L_DIV)
    for ph in list(s.placeholders):
        if str(ph.placeholder_format.type).startswith(("TITLE", "BODY", "OBJECT", "SUBTITLE")):
            ph._element.getparent().remove(ph._element)
    circulo(s, M, 2.55, 1.05, BRANCO)
    txt(s, letra, M, 2.55, 1.05, 1.05, size=34, font=FT, bold=True, cor=AZUL_E,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, titulo_txt, M + 1.5, 2.62, W - M - 2.2, 0.8, size=32, font=FT, bold=True, cor=BRANCO)
    txt(s, sub, M + 1.5, 3.5, W - M - 2.2, 0.5, size=14.5, cor="D6E9F7", italic=True)
    return s

# =====================================================================
# 1 CAPA  (slide padrao do template, apenas preenchido)
SEP = chr(10)

def preencher_ph(slide, idx, texto, size, font, bold=False, cor=None):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            linhas = texto.split(SEP)
            for i, ln in enumerate(linhas):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                r = p.add_run(); r.text = ln
                r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
                if cor: r.font.color.rgb = rgb(cor)
            return True
    return False

preencher_ph(CAPA, 0,  "Como eu uso IA" + SEP + "no dia a dia", 34, FT, True, BRANCO)
preencher_ph(CAPA, 11, "Práticas que funcionam — contadas pelo" + SEP +
                       "caminho que eu fiz até elas", 13, FB, False, BRANCO)
notas(CAPA, "Enquadramento: isto nao e treinamento de ferramenta. Metade do tempo e de voces: 30 min de fala, 30 de conversa.")

# 2 O QUE E / NAO E
s = novo(L_TITULO, "O que esta conversa é")
bloco(s, M, 1.85, 5.8, 2.0, "✕", "O que NÃO é",
      "Treinamento de ferramenta.\nLista de prompts mágicos.\nAlguém dizendo que vocês estão fazendo errado.",
      ERRO, ERR_BG, BORDA_E, size=12.5)
bloco(s, M + 6.1, 1.85, 5.8, 2.0, "✓", "O que é",
      "Práticas que funcionam, com exemplo real.\nO caminho até elas — que no meu caso foi caro.\nO que dá para aplicar já na segunda.",
      VERDE, OK_BG, BORDA_O, size=12.5)
txt(s, "Se no fim vocês discordarem, ótimo — é para isso que metade do tempo está reservada.",
    M, 4.15, W - 2 * M, 0.5, size=15, italic=True)
notas(s, "Baixar a guarda da plateia. Deixar claro que e experiencia, nao auditoria.")

# 3 TRES ERROS
s = novo(L_TITULO, "Três erros que custam caro",
         "Os três eram meus. Pela conversa com outras squads, não eram só meus.")
itens = [("Colar tudo", "Log inteiro, arquivo inteiro, “por precaução”.\nA intuição de que mais contexto = melhor resposta."),
         ("Pedir vago", "“Otimiza isso aí.”\nVolta algo legítimo que não dá para usar."),
         ("Aceitar o que vem", "Bem escrito, parece certo, ninguém roda.\nAparece no review — ou em produção.")]
y = 2.05
for i, (t1, t2) in enumerate(itens, 1):
    circulo(s, M, y, 0.46, ERRO)
    txt(s, str(i), M, y, 0.46, 0.46, size=13, bold=True, cor=BRANCO,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, t1, M + 0.72, y - 0.02, 3.3, 0.4, size=16, bold=True)
    txt(s, t2, M + 4.1, y - 0.04, 7.4, 0.8, size=12.5, cor=MUTED, space=1.25)
    y += 1.18
rodape(s, "Nenhum dos três é sobre a ferramenta. Os três são sobre método.")
notas(s, "Contar um caso real curto de um dos tres. Vulnerabilidade primeiro.")

# 4 GARGALO
s = novo(L_TITULO, "O gargalo mudou de lugar")
txt(s, "Escrever código deixou de ser o gargalo.", M, 2.2, 11.8, 0.6, size=25, font=FT)
txt(s, "Decidir o que deve ser construído virou o gargalo.", M, 2.95, 11.8, 0.6,
    size=25, font=FT, bold=True, cor=AZUL)
txt(s, "Quando ninguém decide, o modelo decide — por omissão, com confiança.\nE isso só aparece no code review.",
    M, 4.2, 11.4, 0.9, size=15, cor=MUTED, space=1.35)
notas(s, "Slide central. Se levarem uma frase, e esta.")

# 5 CICLO
s = novo(L_TITULO, "O ciclo",
         "O mesmo, de tarefa pequena a feature grande — só muda o rigor de cada etapa")
passos = [("ENTENDER", "ler antes\nde pedir"), ("ESPECIFICAR", "escrever o\ncritério"),
          ("EXECUTAR", "pessoa ou\nagente"), ("VERIFICAR", "rodar, não\nler"),
          ("REGISTRAR", "o que vale\npra próxima")]
bw, gap, ytop = 2.12, 0.34, 2.35
x = (W - (len(passos) * bw + (len(passos) - 1) * gap)) / 2
for i, (a, b) in enumerate(passos):
    dest = (i == 1)
    card(s, x, ytop, bw, 1.7, AZUL if dest else SUAVE, AZUL if dest else None)
    txt(s, a, x, ytop + 0.28, bw, 0.34, size=12, bold=True,
        cor=BRANCO if dest else TEXTO, align=PP_ALIGN.CENTER)
    txt(s, b, x, ytop + 0.72, bw, 0.7, size=11, cor=BRANCO if dest else MUTED,
        align=PP_ALIGN.CENTER, space=1.2)
    if i < len(passos) - 1:
        txt(s, "→", x + bw, ytop + 0.6, gap, 0.45, size=15, cor=CINZA,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += bw + gap
txt(s, "A etapa 2 é a que mais se pula — e a única que, quando falha, estraga todas as outras.",
    M, 4.5, W - 2 * M, 0.5, size=15, italic=True)
rodape(s, "O atalho comum — e o meu, por muito tempo: pular de “tenho um problema” direto para “escreve o código”.")
notas(s, "Mostrar que o ciclo nao e cerimonia: em tarefa pequena cada etapa leva segundos.")

# 6 DIVISOR A
notas(divisor("A", "Antes de pedir", "A parte que acontece antes de digitar qualquer coisa"),
      "Transicao curta.")

# 7 CONTEXTO
s = novo(L_TITULO, "Contexto é orçamento, não depósito")
bloco(s, M, 1.72, 5.8, 3.0, "✕", "Sem recorte",
      "Aqui o log, me ajuda:\n\n09:14:02 INFO  Starting worker\n09:14:02 INFO  Connected redis\n  ... (1.847 linhas) ...\n09:31:55 ERROR ValueError: inva\n  ... (mais 600 linhas) ...",
      ERRO, ERR_BG, BORDA_E, mono=True, size=10.5)
bloco(s, M + 6.1, 1.72, 5.8, 3.0, "✓", "Com recorte",
      "Erro no worker, ~2% dos registros.\n\nValueError: invalid literal\n  parser.py:42 in parse_row\n\nJá verifiquei: CSV tem células\nvazias. Schema não é meu.\n\nQuero vazio -> None, sem engolir\noutros erros.",
      VERDE, OK_BG, BORDA_O, mono=True, size=10.5)
txt(s, "~40.000 tokens", M, 4.88, 5.8, 0.45, size=18, bold=True, cor=ERRO, align=PP_ALIGN.CENTER)
txt(s, "~200 tokens · e resposta melhor", M + 6.1, 4.88, 5.8, 0.45, size=18, bold=True,
    cor=AZUL_E, align=PP_ALIGN.CENTER)
notas(s, "O ponto nao e o custo: as 1.847 linhas irrelevantes disputam atencao com as 6 que importam.")

# 8 SESSAO LONGA
s = novo(L_TITULO, "Sessão longa apodrece",
         "Conversa de 3 horas acumula decisão revogada, tentativa abandonada e arquivo velho")
card(s, M, 2.1, W - 2 * M, 1.9, SUAVE)
txt(s, "O que pedir antes de abrir uma sessão nova:", M + 0.32, 2.3, 10, 0.34, size=13, bold=True)
txt(s, "“Escreva um resumo de estado em até 20 linhas: decisões tomadas e o porquê,\ndecisões que DESCARTAMOS e o motivo, arquivos alterados, o que está pendente.”",
    M + 0.32, 2.74, 11.4, 1.0, size=12, font=FM, space=1.3)
txt(s, "~90.000 tokens de histórico  →  ~600 de estado útil", M, 4.3, W - 2 * M, 0.5,
    size=18, bold=True, cor=AZUL)
txt(s, "O bloco de decisões descartadas é o que mais economiza: sem ele, a sessão nova propõe de novo tudo que já tinha sido recusado.",
    M, 4.95, 11.5, 0.7, size=13.5, cor=MUTED, space=1.3)
notas(s, "Sugerir que testem hoje numa conversa longa que ja esteja aberta.")

# 9 DIVISOR B
notas(divisor("B", "Como pedir", "Onde o retorno aparece mais rápido"), "Transicao curta.")

# 10 CRITERIO DE ACEITE
s = novo(L_TITULO, "O critério de aceite vai junto com o pedido")
bloco(s, M, 1.72, 5.8, 3.1, "✕", "Pedido vago",
      "otimiza essa função aqui pra mim", ERRO, ERR_BG, BORDA_E, mono=True, size=11)
txt(s, "Otimizar em quê? Tempo, memória, número de queries?\n\nO modelo escolhe — e devolve algo legítimo, com async, ORM trocado e assinatura nova.\n\nNão dá pra usar.",
    M + 0.26, 2.62, 5.3, 1.9, size=12, cor=TEXTO, space=1.25)
bloco(s, M + 6.1, 1.72, 5.8, 3.1, "✓", "Pedido com aceite",
      "Suspeita: N+1 no lazy load.\n\nRESTRIÇÕES\n- manter assinatura e retorno\n- sem dependência nova\n- sem async\n\nACEITE\n- nº de queries constante\n- diga quantas antes e depois",
      VERDE, OK_BG, BORDA_O, mono=True, size=10.5)
txt(s, "O pedido errado devolve código. O certo devolve código que entra no PR.",
    M, 5.0, W - 2 * M, 0.5, size=15, bold=True)
notas(s, "Custo de escrever: 40 segundos. Custo de nao escrever: o review inteiro.")

# 11 PLANO ANTES
s = novo(L_TITULO, "Em tarefa grande: plano antes do código",
         "O parágrafo que mais poupa retrabalho")
card(s, M, 2.1, 7.3, 2.25, NAVY if not DARK else "16203A")
txt(s, "Antes de escrever qualquer código:\n1. liste os arquivos que pretende tocar\n2. aponte os riscos e o que pode quebrar\n3. diga o que você NÃO vai mexer\n\nPare aí e espere meu OK.",
    M + 0.32, 2.34, 6.6, 1.85, size=12, font=FM, cor=BRANCO, space=1.3)
txt(s, "Por que funciona", M + 7.8, 2.14, 4.0, 0.4, size=15, bold=True)
txt(s, "•  Corrigir um plano de 10 linhas custa uma frase.\n\n•  Corrigir 400 linhas geradas custa a tarefa inteira.\n\n•  “Não mexe no serializers.py, é contrato público” — o modelo não tinha como saber.",
    M + 7.8, 2.64, 4.05, 1.9, size=12, cor=MUTED, space=1.2)
rodape(s, "Vale sempre que a tarefa passa de um arquivo.")
notas(s, "Demonstrar ao vivo se der tempo: e o item mais facil de copiar.")

# 12 DIVISOR C
notas(divisor("C", "O que persistir", "Contexto morre com a sessão. O que precisa sobreviver vira arquivo"),
      "Transicao curta.")

# 13 ARQUIVO DE CONTEXTO
s = novo(L_TITULO, "Um arquivo de contexto no repositório",
         "Lido em toda sessão, por todo mundo — então cada linha precisa se pagar")
bloco(s, M, 2.05, 5.8, 2.55, "✕", "Não entra",
      "“Usamos Python e FastAPI.”\n“Os testes ficam em tests/.”\n“Escreva código de qualidade.”\n\nO modelo descobre isso em 2 segundos lendo o repo. E o time paga por essas linhas todo dia.",
      ERRO, ERR_BG, BORDA_E, size=12)
bloco(s, M + 6.1, 2.05, 5.8, 2.55, "✓", "Entra",
      "“app/legacy/ está congelado.”\n“Nunca chame billing direto — o rate limit é compartilhado e derruba geral.”\n“PA = Pedido Antecipado.”\n\nIsso nenhum modelo descobre sozinho.",
      VERDE, OK_BG, BORDA_O, size=12)
txt(s, "Teste: se a linha pode ser respondida lendo o README, apague.",
    M, 4.85, W - 2 * M, 0.5, size=15, bold=True)
notas(s, "Armadilhas e a secao de maior valor. Pergunta: o que ja quebrou alguem que entrou agora?")

# 14 MEMORIA
s = novo(L_TITULO, "Onde mora o que o código não conta",
         "Vault de notas, wiki, ADR no repo — a ferramenta importa menos que o formato")
bloco(s, M, 2.0, 5.8, 2.8, "✕", "Nota-diário (inútil)",
      "“Ontem discutimos o importador.\nO pessoal falou que era melhor a\nabordagem nova. O João vai olhar.”\n\nSeis meses depois: qual problema?\nqual abordagem? “ontem” quando?",
      ERRO, ERR_BG, BORDA_E, size=11.5)
bloco(s, M + 6.1, 2.0, 5.8, 2.8, "✓", "Fato, com porquê e validade",
      "“Importação roda em lote noturno.\n\nPor quê: fornecedor limita 100 req/min.\nTempo real derrubava a integração.\n\nConsequência: estoque tem 24h de atraso.\n\nDecidido em 2026-05-20.”",
      VERDE, OK_BG, BORDA_O, size=11.5)
txt(s, "O título é a afirmação, não o tópico — é por ele que a IA decide se a nota importa.",
    M, 5.0, W - 2 * M, 0.5, size=14.5, italic=True)
notas(s, "Regra: data absoluta sempre. 'Recentemente' apodrece em silencio.")

# 15 METODO NA FERRAMENTA
s = novo(L_TITULO, "Método dentro da ferramenta",
         "O que mais muda resultado não é trocar de modelo — é o ciclo deixar de ser opcional")
pares = [("Com framework", "O ciclo é imposto: entender → planejar → executar → verificar.\nO agente não avança de etapa sem cumprir a anterior.", True),
         ("Sem framework", "O ciclo é seguido quando dá tempo.\nSob prazo, é a primeira coisa a ser abandonada.", False)]
y = 2.15
for nome, desc, dest in pares:
    card(s, M, y, W - 2 * M, 1.2, AZUL if dest else SUAVE, AZUL if dest else None)
    txt(s, nome, M + 0.32, y + 0.2, 3.0, 0.4, size=15, bold=True, cor=BRANCO if dest else TEXTO)
    txt(s, desc, M + 3.5, y + 0.2, 8.0, 0.85, size=12.5,
        cor=BRANCO if dest else MUTED, space=1.25)
    y += 1.42
txt(s, "Disciplina que depende de alguém lembrar, sob prazo, não é disciplina.",
    M, 5.15, 11.8, 0.5, size=17, font=FT, bold=True)
notas(s, "Dizer qual voce usa (Superpowers) como exemplo, nao como recomendacao fechada. O processo custa tokens a mais por sessao; o ganho aparece em nao refazer.")

# 16 DIVISOR D
notas(divisor("D", "Quando delegar", "O critério vale para qualquer agente autônomo"),
      "Se a squad usa Devin, abrir a conversa aqui.")

# 17 MATRIZ
s = novo(L_TITULO, "Na mão · piloto · delega",
         "Decisão que vem antes de abrir a ferramenta, não depois")
cols = [("NA MÃO", ERRO, ERR_BG, "Decisão de arquitetura\nTrade-off de produto\nCódigo de auth e pagamento\nIncidente em produção"),
        ("PILOTO", AZUL, NEU_BG, "Causa desconhecida\nDesign ainda aberto\nRefactor que atravessa módulos\nExploração de alternativas"),
        ("DELEGA", VERDE, OK_BG, "Migração mecânica\nBug com repro determinístico\nTeste faltando\nLint e tipo em massa")]
bw, gap = 3.82, 0.33
x = M
for nome, cor, bg, corpo in cols:
    card(s, x, 2.05, bw, 2.7, bg)
    txt(s, nome, x + 0.24, 2.24, bw - 0.48, 0.4, size=14, bold=True, cor=cor)
    txt(s, corpo, x + 0.24, 2.72, bw - 0.48, 1.9, size=12, space=1.45)
    x += bw + gap
txt(s, "Delegue o tedioso e verificável. Pilote o ambíguo. Faça na mão o irreversível.",
    M, 4.95, W - 2 * M, 0.45, size=15, bold=True)
txt(s, "A pergunta que resolve quase tudo: existe teste que prova que ficou certo? Se não existe, não delegue.",
    M, 5.42, 11.8, 0.45, size=13, cor=MUTED, italic=True)
notas(s, "Perguntar em qual coluna cai o que eles delegam hoje.")

# 18 REVIEW
s = novo(L_TITULO, "Um PR que passa no CI e não deve ser aprovado")
card(s, M, 1.72, 6.5, 2.9, NAVY if not DARK else "16203A")
txt(s, "  def parse_row(row) -> int | None:\n-     return int(row[\"quantidade\"])\n+     try:\n+         return int(row[\"quantidade\"])\n+     except Exception:\n+         return None\n\n- def test_invalido():\n-     with pytest.raises(ValueError):\n+ def test_invalido():\n+     assert parse_row(\"abc\") is None",
    M + 0.28, 1.95, 5.9, 2.45, size=10, font=FM, cor=BRANCO, space=1.15)
txt(s, "CI verde. Três problemas:", M + 7.0, 1.76, 4.9, 0.4, size=15, bold=True)
txt(s, "•  except Exception engole tudo — coluna faltando vira None silencioso\n\n•  o teste antigo foi invertido para passar\n\n•  e tem um parâmetro de biblioteca que simplesmente não existe",
    M + 7.0, 2.3, 4.9, 2.2, size=12, cor=MUTED, space=1.2)
txt(s, "Ler o diff não é revisar. O segundo item é o mais grave e o mais fácil de passar batido.",
    M, 4.8, W - 2 * M, 0.45, size=15, bold=True)
txt(s, "Atalho: git diff nos testes primeiro. Linha alterada em teste existente exige justificativa escrita.",
    M, 5.26, 11.8, 0.45, size=13, cor=MUTED, italic=True)
notas(s, "Deixar a plateia achar os problemas antes de revelar.")

# 19 ONDE NAO USAR
s = novo(L_TITULO, "Onde não usar IA", "Saber onde não usar é metade do ganho")
itens = [("Conta e agregação", "O modelo estima a partir do texto. Peça a query — a máquina calcula, o número é auditável."),
         ("Decisão de arquitetura", "Serve para levantar alternativas. A escolha é da pessoa — e é onde está o valor dela."),
         ("Aprovar o que não se entende", "Se não entende cada linha, não aprove. Dívida no dia 1, incidente no dia 90."),
         ("Dado de cliente e segredo", "Nunca no prompt. Anonimize ou use dado sintético. Sem exceção de prazo.")]
y = 2.05
for t1, t2 in itens:
    circulo(s, M, y, 0.38, ERRO)
    txt(s, "✕", M, y, 0.38, 0.38, size=11, bold=True, cor=BRANCO,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, t1, M + 0.6, y - 0.03, 3.6, 0.4, size=14, bold=True)
    txt(s, t2, M + 4.3, y - 0.03, 7.3, 0.7, size=12, cor=MUTED, space=1.2)
    y += 0.88
rodape(s, "A responsabilidade não é delegável: quem aprova responde. “A IA escreveu” não explica incidente.")
notas(s, "O item de dado de cliente merece pausa se a squad lida com dado regulado.")

# 20 FERRAMENTAS
s = novo(L_TITULO, "Ferramentas, rápido",
         "O que eu uso hoje. Menos importante que tudo que veio antes — mas sempre perguntam")
linhas = [("Claude Code", "onde o trabalho acontece"),
          ("Superpowers", "framework de processo: impõe o ciclo dentro da ferramenta"),
          ("Obsidian", "o que o código não conta: decisão, porquê, regra de negócio"),
          ("Compressão de saída", "resposta densa; economiza agora e no histórico depois")]
y = 2.1
for a, b in linhas:
    card(s, M, y, W - 2 * M, 0.78, SUAVE)
    txt(s, a, M + 0.28, y + 0.2, 3.2, 0.4, size=14, bold=True, cor=AZUL)
    txt(s, b, M + 3.7, y + 0.2, 7.9, 0.4, size=12.5)
    y += 0.92
txt(s, "Trocar de ferramenta não resolve problema de método. Pedido ruim é ruim em qualquer uma.",
    M, 5.85, W - 2 * M, 0.45, size=15, bold=True)
notas(s, "Passar rapido. Se perguntarem de Devin: nao uso no dia a dia, mas o criterio da matriz vale igual.")

# 21 SEGUNDA-FEIRA
s = novo(L_TITULO, "Três coisas para segunda-feira",
         "Nessa ordem — a primeira sozinha já muda o resultado")
itens = [("1", "Mande o trecho, não o arquivo", "Recorte o traceback e diga o que já verificou."),
         ("2", "Diga como você vai saber que ficou certo", "Uma linha de critério de aceite, junto com o pedido."),
         ("3", "Rode antes de aprovar", "Ler o diff não é revisar. Vale para PR de IA e de gente.")]
y = 2.15
for num, t1, t2 in itens:
    card(s, M, y, W - 2 * M, 1.1, SUAVE)
    circulo(s, M + 0.3, y + 0.28, 0.52, AZUL)
    txt(s, num, M + 0.3, y + 0.28, 0.52, 0.52, size=14, bold=True, cor=BRANCO,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, t1, M + 1.1, y + 0.22, 10.3, 0.4, size=15, bold=True)
    txt(s, t2, M + 1.1, y + 0.62, 10.3, 0.35, size=12, cor=MUTED)
    y += 1.26
rodape(s, "Nenhuma das três depende de instalar nada.")
notas(s, "Fechar pedindo que escolham UMA e tentem essa semana.")

# 22 Q&A  (slide de conteudo: nao pode dividir espaco com o logo do fechamento)
s = novo(L_TITULO, "Agora é de vocês")
txt(s, "30 minutos. Discordar é o melhor uso desse tempo.", M, 1.45, 11, 0.5,
    size=16, bold=True, cor=AZUL_E)
itens_qa = [("Onde isso não vai funcionar no nosso contexto?", ""),
            ("O que vocês já fazem que o resto do time deveria copiar?", ""),
            ("Qual tarefa da sprint serviria de teste?", "a mais concreta — começa por essa se travar")]
y = 2.3
for pergunta, nota in itens_qa:
    card(s, M, y, W - 2 * M, 0.95, SUAVE)
    txt(s, pergunta, M + 0.32, y + (0.3 if not nota else 0.16), 10.8, 0.4, size=15, bold=True)
    if nota:
        txt(s, nota, M + 0.32, y + 0.56, 10.8, 0.3, size=12, cor=MUTED, italic=True)
    y += 1.12
txt(s, "Material completo e templates: <link do repositório>", M, 5.95, 10, 0.4,
    size=12, cor=MUTED, italic=True)
notas(s, "Se travar, comecar pela terceira pergunta: e a mais concreta e sempre destrava.")

# 23 ENCERRAMENTO NTT  (layout oficial, sem nada por cima do logo)
s = prs.slides.add_slide(L_FIM)
for ph in list(s.placeholders):
    if str(ph.placeholder_format.type).startswith(("TITLE", "BODY", "OBJECT", "SUBTITLE")):
        ph._element.getparent().remove(ph._element)
notas(s, "Slide de encerramento da marca. Nada e sobreposto ao logo.")

prs.save(SAIDA)
print("ok: %s  (%d slides, estilo=%s)" % (SAIDA, len(prs.slides._sldIdLst), ESTILO))
